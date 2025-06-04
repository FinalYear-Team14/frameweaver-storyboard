"""
Enhanced storyboard generator with improved character consistency, full-color comic generation, and context tracking
"""

import os
import time
import json
import logging
import mimetypes
import pymysql
from dotenv import load_dotenv
from PIL import Image, ImageDraw, ImageFont
from fpdf import FPDF
import spacy
from openai import OpenAI
from datetime import datetime
from tabulate import tabulate
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import requests
import hashlib
import base64

# -------------------
# 1) ENV + LOGGING
# -------------------
load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.FileHandler("app.log"), logging.StreamHandler()],
)

# --------------------------
# 2) CONFIGURE API CLIENTS
# --------------------------
nlp = spacy.load("en_core_web_sm")

openai_api_key = os.getenv("OPENAI_API_KEY")

# MySQL Config
mysql_host = os.getenv("MYSQL_HOST", "localhost")
mysql_user = os.getenv("MYSQL_USER", "root")
mysql_password = os.getenv("MYSQL_PASSWORD", "")
mysql_db = os.getenv("MYSQL_DB", "story_database")

# Initialize OpenAI client
if not openai_api_key or openai_api_key.strip() == "":
    logging.error("No valid OpenAI API key found. Text enhancement, image generation, and audio generation will fail.")
    openai_client = None
else:
    openai_client = OpenAI(api_key=openai_api_key)
    logging.info("OpenAI client initialized successfully")


# --------------------------
# CONTEXT TRACKING CLASS
# --------------------------
class PanelContextTracker:
    """Tracks details from previous panels to maintain consistency."""

    def __init__(self):
        self.panels = []
        self.character_last_seen = {}
        self.character_details = {}
        self.character_poses = {}
        self.character_positions = {}
        self.scene_settings = {}

    def add_panel(self, panel_number, scene_text, characters_present, character_descriptions, setting_details=None):
        """Add a panel's details to the context."""
        panel_info = {
            'panel_number': panel_number,
            'scene_text': scene_text,
            'characters_present': characters_present,
            'character_descriptions': character_descriptions,
            'setting_details': setting_details,
            'timestamp': datetime.now()
        }
        self.panels.append(panel_info)

        # Update last seen info for each character
        for char_name in characters_present:
            self.character_last_seen[char_name] = panel_number
            if char_name in character_descriptions:
                self.character_details[char_name] = character_descriptions[char_name]

        # Store setting details
        if setting_details:
            self.scene_settings[panel_number] = setting_details

    def get_context_for_panel(self, current_panel, characters_in_scene, scene_text=""):
        """Get context from previous panels for the current panel."""
        context_parts = []

        # Get last 3 panels for reference
        recent_panels = self.panels[-3:] if len(self.panels) > 0 else []

        if recent_panels:
            context_parts.append("CONTEXT FROM PREVIOUS PANELS:")

            # Describe recent panels
            for panel in recent_panels:
                panel_desc = f"Panel {panel['panel_number']}: {', '.join(panel['characters_present'])} appeared"
                if panel.get('setting_details'):
                    panel_desc += f" in {panel['setting_details']}"
                context_parts.append(f"- {panel_desc}")

        # Add specific character continuity info
        context_parts.append("\nCHARACTER CONTINUITY:")
        for char_name in characters_in_scene:
            char_context = []

            if char_name in self.character_last_seen:
                last_panel = self.character_last_seen[char_name]
                if last_panel == current_panel - 1:
                    char_context.append(f"{char_name} was just shown in panel {last_panel}")
                    char_context.append("MUST maintain exact same appearance as previous panel")
                else:
                    char_context.append(f"{char_name} was last seen in panel {last_panel}")
                    char_context.append(f"Ensure {char_name} looks IDENTICAL to their appearance in panel {last_panel}")

                if char_name in self.character_details:
                    details = self.character_details[char_name]
                    char_context.append(f"Previous appearance details: {details}")

            if char_context:
                context_parts.append(f"- {char_name}: {' | '.join(char_context)}")

        # Add setting continuity if same location
        if recent_panels and self._is_same_location(scene_text, recent_panels[-1].get('scene_text', '')):
            context_parts.append("\nSETTING CONTINUITY:")
            context_parts.append("- This scene continues in the SAME LOCATION as the previous panel")
            context_parts.append("- Maintain consistent background, lighting, and atmosphere")

        return "\n".join(context_parts) if context_parts else ""

    def _is_same_location(self, current_scene, previous_scene):
        """Determine if two scenes are in the same location."""
        # Simple heuristic - look for location keywords
        location_keywords = ['room', 'street', 'house', 'office', 'park', 'building', 'car', 'forest', 'city']

        current_locations = [word for word in current_scene.lower().split() if word in location_keywords]
        previous_locations = [word for word in previous_scene.lower().split() if word in location_keywords]

        return bool(set(current_locations) & set(previous_locations))


# Fallback image generation function if OpenAI auth fails
def fallback_image_generation(prompt, output_path):
    """Creates a simple placeholder image with text and color."""
    try:
        width, height = 800, 600
        # Use a light background color instead of white
        image = Image.new('RGB', (width, height), color=(240, 248, 255))  # Alice Blue
        draw = ImageDraw.Draw(image)

        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except:
            font = ImageFont.load_default()

        # Colored border
        draw.rectangle([(10, 10), (width - 10, height - 10)], outline=(70, 130, 180), width=3)  # Steel Blue

        title = "Image Generation Failed - Placeholder Image"
        draw.text((width // 2, 50), title, fill=(220, 20, 60), font=font, anchor="mm")  # Crimson

        y_pos = 100
        words = prompt.split()
        line = ""
        for word in words:
            test_line = line + word + " "
            if len(test_line) * 10 < width - 40:
                line = test_line
            else:
                draw.text((20, y_pos), line, fill=(25, 25, 112), font=font)  # Midnight Blue
                y_pos += 30
                line = word + " "

            if y_pos > height - 50:
                break

        if line:
            draw.text((20, y_pos), line, fill=(25, 25, 112), font=font)

        image.save(output_path)
        return True
    except Exception as e:
        logging.error(f"Failed to create fallback image: {e}")
        return False


# ----------------------------
# 3) DATABASE INITIALIZATION
# ----------------------------
def initialize_database():
    """Create database and tables if they don't exist."""
    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password
        )
        cursor = conn.cursor()

        cursor.execute(f"CREATE DATABASE IF NOT EXISTS {mysql_db}")
        cursor.execute(f"USE {mysql_db}")

        cursor.execute("""
            CREATE TABLE IF NOT EXISTS story_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                title VARCHAR(255) NOT NULL,
                story_outline TEXT,
                enhanced_story TEXT,
                audio_urls JSON,
                image_urls JSON,
                pdf_url VARCHAR(255),
                pptx_url VARCHAR(255),
                user_id INT,
                character_styles JSON,
                character_reference_sheets JSON,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        cursor.execute("""
            CREATE TABLE IF NOT EXISTS users (
                id INT AUTO_INCREMENT PRIMARY KEY,
                username VARCHAR(50) UNIQUE NOT NULL,
                email VARCHAR(100) UNIQUE NOT NULL,
                password VARCHAR(255) NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        conn.commit()
        cursor.close()
        conn.close()
        logging.info("Database initialization successful")
    except Exception as e:
        logging.error(f"Error initializing database: {e}")


# -----------------------
# 4) TEXT ENHANCEMENT
# -----------------------
def text_enhancer(story_outline):
    """Enhances a story outline using OpenAI GPT-based model."""
    logging.info("Starting text enhancement...")
    start_time = time.time()

    if not openai_client:
        logging.error("OpenAI client not initialized. Cannot enhance story.")
        return story_outline

    messages = [{"role": "system", "content": f"Enhance this story: {story_outline}"}]

    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o", messages=messages, temperature=0.5, max_tokens=2048
        )
        enhanced_story = response.choices[0].message.content.strip()
        logging.info(f"Text enhancement completed in {time.time() - start_time:.2f} seconds.")
        return enhanced_story
    except Exception as e:
        logging.error(f"Error in enhancing story: {e}")
        return story_outline


# ------------------------------------------
# 5) SAVE STORY LOCALLY (TXT + JSON)
# ------------------------------------------
def save_story_to_file(story_text, title, output_dir="stories"):
    """Saves the story to text and JSON files within a dedicated directory."""
    logging.info(f"Saving story '{title}' to disk...")

    os.makedirs(output_dir, exist_ok=True)
    safe_title = "".join(c if c.isalnum() else "_" for c in title)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{safe_title}_{timestamp}"

    txt_path = os.path.join(output_dir, f"{filename}.txt")
    json_path = os.path.join(output_dir, f"{filename}.json")

    try:
        with open(txt_path, "w", encoding="utf-8") as txt_file:
            txt_file.write(story_text)
        with open(json_path, "w", encoding="utf-8") as json_file:
            json.dump({"title": title, "story": story_text, "timestamp": timestamp}, json_file, indent=4)
        logging.info(f"Story saved to {txt_path} and {json_path}")
        return txt_path, json_path
    except Exception as e:
        logging.error(f"Error saving story: {e}")
        return None, None


# -------------------------------------------------------
# 6) TEXT SANITIZATION FOR PDF
# -------------------------------------------------------
def sanitize_text_for_pdf(text):
    """Sanitize text to be compatible with FPDF latin-1 encoding."""
    replacements = {
        '\u2014': '-',
        '\u2013': '-',
        '\u2018': "'",
        '\u2019': "'",
        '\u201c': '"',
        '\u201d': '"',
        '\u2026': '...',
        '\u2022': '*',
        '\u00a0': ' ',
    }

    for unicode_char, ascii_char in replacements.items():
        text = text.replace(unicode_char, ascii_char)

    return re.sub(r'[^\x00-\xFF]', '?', text)


# -------------------------------------------------------
# 7) ENHANCED CHARACTER EXTRACTION WITH COLOR
# -------------------------------------------------------
def extract_characters_from_story(story_text):
    """Extract character names and ultra-detailed descriptions with color information."""
    logging.info("Extracting character information from story...")

    if not openai_client:
        logging.error("OpenAI client not initialized. Cannot extract characters.")
        return {}

    try:
        messages = [
            {"role": "system",
             "content": "You are a character design specialist who extracts extremely precise visual specifications including colors."},
            {"role": "user", "content": f"""
                Extract the main characters from this story. For each character, provide ULTRA-SPECIFIC details with COLOR INFORMATION:

                1. Name
                2. Exact height in feet/inches (e.g., "5'7")
                3. Precise age (e.g., "32 years old")
                4. Face shape using geometric terms (e.g., "oval face with 1.4:1 length-to-width ratio")
                5. Hair: exact color (with hex code if applicable), length in inches, style (e.g., "#8B4513 chestnut brown, 14-inch length, straight with slight wave")
                6. Eyes: shape, color (with hex code), spacing (e.g., "almond-shaped #00CED1 turquoise eyes, 1.5 eye-widths apart")
                7. Nose: shape and size (e.g., "straight Roman nose, 2.1 inches long")
                8. Skin tone: specific shade with color reference (e.g., "Fitzpatrick Type III, warm beige #D2B48C")
                9. Build: specific proportions (e.g., "athletic build, 38-30-38 measurements, BMI 22")
                10. Clothing: exact items with specific colors (e.g., "royal blue blazer #4169E1, crisp white oxford shirt #FFFFFF, charcoal grey pants #36454F")
                11. Distinctive features: precise descriptions (e.g., "thin scar 1.2 inches long on left eyebrow")
                12. Typical pose and expression (e.g., "slightly tilted head, gentle smile showing upper teeth")

                Format the response as a JSON object with character names as keys.
                Story:
                {story_text}
            """}
        ]

        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            temperature=0.1,
            max_tokens=3000,
            response_format={"type": "json_object"}
        )

        characters_data = json.loads(response.choices[0].message.content)
        logging.info(f"Extracted {len(characters_data)} characters from story")
        return characters_data

    except Exception as e:
        logging.error(f"Error extracting characters from story: {e}")
        return {}


# -------------------------------------------------------
# 8) CHARACTER REFERENCE SHEET GENERATION WITH COLOR
# -------------------------------------------------------
def generate_character_reference_sheet(character_name, character_info, story_dir):
    """Generate a detailed character reference sheet in full color."""
    logging.info(f"Generating character reference sheet for {character_name}")

    if not openai_client:
        return None

    try:
        reference_path = os.path.join(story_dir, f"reference_{character_name.lower().replace(' ', '_')}.png")

        # Create ultra-detailed prompt with COLOR
        prompt = f"""Create a professional character reference sheet for comic/manga production in FULL COLOR:

CHARACTER: {character_name}

PHYSICAL SPECIFICATIONS:
- Height: {character_info.get('height', '5.8')}
- Age: {character_info.get('age', '25')}
- Face Shape: {character_info.get('face_shape', 'oval')}
- Hair: {character_info.get('hair', 'brown, shoulder-length')}
- Eyes: {character_info.get('eyes', 'brown, almond-shaped')}
- Skin Tone: {character_info.get('skin_tone', 'fair')}
- Build: {character_info.get('build', 'average')}

CLOTHING/STYLE:
{character_info.get('clothing', 'casual modern attire')}

DISTINCTIVE FEATURES:
{character_info.get('distinctive_features', 'none')}

REFERENCE SHEET REQUIREMENTS:
- Front view (full body)
- Side profile (head and shoulders)
- 3/4 view (full body)
- Facial expressions chart (happy, sad, angry, surprised)
- Clothing details close-up
- Color reference boxes
- Clear labels for each element
- Professional character sheet layout
- FULL COLOR comic art style
- Include color swatches for hair, eyes, skin, and clothing

STYLE: Clean line art with vibrant colors, professional comic book character sheet"""

        response = openai_client.images.generate(
            model="gpt-image-1",
            prompt=prompt,
            size="1024x1024",
            quality="high",
            n=1,
        )

        if hasattr(response, 'data') and response.data:
            data_item = response.data[0]
            image_url = None

            if hasattr(data_item, 'url') and data_item.url:
                image_url = data_item.url
            elif hasattr(data_item, 'b64_json') and data_item.b64_json:
                img_data = base64.b64decode(data_item.b64_json)
                with open(reference_path, 'wb') as img_file:
                    img_file.write(img_data)
                return reference_path

            if image_url:
                image_response = requests.get(image_url)
                if image_response.status_code == 200:
                    with open(reference_path, 'wb') as img_file:
                        img_file.write(image_response.content)
                    return reference_path

        return None

    except Exception as e:
        logging.error(f"Error generating character reference sheet: {e}")
        return None


# -------------------------------------------------------
# 9) CHARACTER STYLE DEFINITIONS
# -------------------------------------------------------
def create_character_style_definition(character_name, character_info, reference_sheet_path=None):
    """Create a detailed style definition for consistent character rendering."""

    style_definition = {
        "name": character_name,
        "visual_id": hashlib.md5(f"{character_name}_{json.dumps(character_info)}".encode()).hexdigest()[:8],
        "physical_traits": {
            "height": character_info.get('height', '5\'8"'),
            "age": character_info.get('age', '25'),
            "face_shape": character_info.get('face_shape', 'oval face with soft jawline'),
            "hair": character_info.get('hair', 'dark brown, shoulder-length, straight with slight wave'),
            "eyes": character_info.get('eyes', 'brown, almond-shaped, medium-sized'),
            "skin_tone": character_info.get('skin_tone', 'light olive, warm undertone'),
            "build": character_info.get('build', 'athletic, medium frame'),
            "distinctive_features": character_info.get('distinctive_features', 'small scar above left eyebrow')
        },
        "clothing": character_info.get('clothing', 'dark blue jacket, white shirt, brown pants'),
        "style_notes": character_info.get('style_notes', 'confident posture, slight smile'),
        "reference_sheet": reference_sheet_path
    }

    return style_definition


# -------------------------------------------------------
# 10) STYLE GUIDE CREATION WITH COLOR
# -------------------------------------------------------
def create_style_guide_prompt(character_definitions):
    """Create a comprehensive style guide for all characters with color consistency."""

    style_guide = """COMIC STYLE GUIDE - MANDATORY CHARACTER CONSISTENCY (FULL COLOR):

CHARACTER VISUAL IDS AND SPECIFICATIONS:
"""

    for char_name, char_def in character_definitions.items():
        style_guide += f"\n{char_name} [VID:{char_def['visual_id']}]:\n"
        style_guide += f"  Physical: {json.dumps(char_def['physical_traits'], indent=2)}\n"
        style_guide += f"  Clothing: {char_def['clothing']}\n"
        style_guide += f"  Style Notes: {char_def['style_notes']}\n"

    style_guide += """
CONSISTENCY RULES:
1. Each character MUST match their VID specifications EXACTLY including colors
2. NO variations in facial features, hair color, or body proportions
3. Clothing colors and styles remain identical unless explicitly changed in scene
4. Maintain consistent art style and color palette across all panels
5. Use reference sheets as absolute canon for colors
6. Lighting and shading must be consistent with scene mood
7. Color temperature should match the emotional tone of each scene
"""

    return style_guide


# -------------------------------------------------------
# 11) SCENE DETECTION
# -------------------------------------------------------
def extract_key_elements(story_text):
    """Splits story_text into scenes using 'Scene' as markers."""
    scenes, current_scene = [], []
    for line in story_text.splitlines():
        line = line.strip()
        if line.lower().startswith("scene") or line.startswith("**Scene"):
            if current_scene:
                scenes.append(" ".join(current_scene))
                current_scene = []
            current_scene.append(line)
        elif line:
            current_scene.append(line)
    if current_scene:
        scenes.append(" ".join(current_scene))

    if not scenes:
        scenes = [story_text]

    return scenes


def detect_characters_in_scene(scene_text, character_names):
    """Detect which characters appear in a given scene."""
    scene_characters = []
    scene_lower = scene_text.lower()

    for char_name in character_names:
        if char_name.lower() in scene_lower:
            scene_characters.append(char_name)

    return scene_characters if scene_characters else list(character_names)[:2]


# -------------------------------------------------------
# 12) ENHANCED PROMPT BUILDING WITH COLOR AND CONTEXT
# -------------------------------------------------------
def build_consistent_scene_prompt_with_context(scene, character_descriptions, style_guide, scene_num, total_scenes,
                                               context):
    """Build a prompt that enforces consistency with color and includes context from previous panels."""

    prompt = f"""Create FULL COLOR comic panel #{scene_num} of {total_scenes}.

{style_guide}

{context}

CHARACTERS IN THIS SCENE:
{chr(10).join(character_descriptions)}

SCENE DESCRIPTION:
{scene}

CRITICAL REQUIREMENTS:
- Match character VIDs exactly as specified
- Maintain absolute color consistency with previous panels
- Professional comic book quality with vibrant colors
- Clean line art with rich, detailed coloring
- Characters must be instantly recognizable
- Follow the style guide specifications precisely
- Use dynamic shading and lighting for depth
- Maintain consistent color palette throughout
- Ensure continuity with previous panels as described in context"""

    return prompt


# -------------------------------------------------------
# 13) ENHANCED CONSISTENCY VERIFICATION
# -------------------------------------------------------
def verify_character_consistency_enhanced(image_path, character_definitions, scene_characters):
    """Enhanced verification using character definitions with color checking."""
    logging.info("Verifying character consistency with enhanced definitions...")

    if not openai_client:
        return True

    try:
        with open(image_path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')

        # Create detailed verification prompt
        verification_prompt = f"""Analyze this comic panel for character consistency.

EXPECTED CHARACTERS:
{json.dumps({char: character_definitions[char] for char in scene_characters}, indent=2)}

VERIFICATION CHECKLIST:
1. Are all expected characters present?
2. Do facial features match specifications?
3. Is hair color/style consistent?
4. Are body proportions correct?
5. Is clothing identical to reference (including colors)?
6. Are distinctive features visible?
7. Do skin tones match specifications?
8. Are eye colors correct?

Respond with JSON: {{"consistent": boolean, "issues": [...], "confidence": 0-1}}"""

        messages = [
            {"role": "system", "content": "You are a professional comic art consistency checker."},
            {"role": "user", "content": [
                {"type": "text", "text": verification_prompt},
                {"type": "image_url", "image_url": {
                    "url": f"data:image/png;base64,{base64_image}"
                }}
            ]}
        ]

        response = openai_client.chat.completions.create(
            model="gpt-4-vision-preview",
            messages=messages,
            temperature=0.1,
            max_tokens=500,
            response_format={"type": "json_object"}
        )

        verification = json.loads(response.choices[0].message.content)
        return verification.get('consistent', True) and verification.get('confidence', 0) > 0.8

    except Exception as e:
        logging.error(f"Error in enhanced consistency verification: {e}")
        return True


# -------------------------------------------------------
# 14) ENHANCED IMAGE GENERATION WITH COLOR CONSISTENCY AND CONTEXT
# -------------------------------------------------------
def generate_comic_images_with_consistency(story_text, title, story_dir, character_styles=None, consistency_level=0.8):
    """Enhanced version with character reference sheets, color style definitions, and context tracking."""
    logging.info("Starting style-consistent color comic image generation with reference sheets and context.")

    if not openai_client:
        logging.error("OpenAI client not initialized. Cannot generate images.")
        return []

    image_dir = os.path.join(story_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    scenes = extract_key_elements(story_text)
    image_paths = []

    # Initialize context tracker
    context_tracker = PanelContextTracker()

    # Extract character info if not provided
    if not character_styles:
        character_styles = extract_characters_from_story(story_text)

    # Generate reference sheets for all characters
    character_definitions = {}
    reference_sheets = {}

    for char_name, char_info in character_styles.items():
        reference_sheet = generate_character_reference_sheet(char_name, char_info, story_dir)
        reference_sheets[char_name] = reference_sheet
        character_definitions[char_name] = create_character_style_definition(
            char_name, char_info, reference_sheet
        )

    # Create a style guide prompt section
    style_guide = create_style_guide_prompt(character_definitions)

    # Generate scene images
    for idx, scene in enumerate(scenes, start=1):
        image_path = os.path.join(image_dir, f"scene_{idx}.png")

        # Detect characters in scene
        scene_characters = detect_characters_in_scene(scene, character_styles.keys())

        # Get context from previous panels
        panel_context = context_tracker.get_context_for_panel(idx, scene_characters, scene)

        # Create character-specific descriptions
        character_descriptions = {}
        character_descriptions_list = []
        for char_name in scene_characters:
            if char_name in character_definitions:
                char_def = character_definitions[char_name]
                char_desc = f"{char_name} [ID:{char_def['visual_id']}]: {json.dumps(char_def['physical_traits'])}"
                character_descriptions_list.append(char_desc)
                character_descriptions[char_name] = char_desc

        # Build the prompt with all consistency elements and context
        prompt = build_consistent_scene_prompt_with_context(
            scene,
            character_descriptions_list,
            style_guide,
            idx,
            len(scenes),
            panel_context
        )

        # Try multiple times for consistency
        max_attempts = 3
        attempt_successful = False

        for attempt in range(max_attempts):
            try:
                response = openai_client.images.generate(
                    model="gpt-image-1",
                    prompt=prompt,
                    size="1024x1024",
                    quality="high",
                    n=1,
                )

                if hasattr(response, 'data') and response.data:
                    data_item = response.data[0]
                    image_url = None

                    if hasattr(data_item, 'url') and data_item.url:
                        image_url = data_item.url
                    elif hasattr(data_item, 'b64_json') and data_item.b64_json:
                        img_data = base64.b64decode(data_item.b64_json)
                        with open(image_path, 'wb') as img_file:
                            img_file.write(img_data)

                        if verify_character_consistency_enhanced(
                                image_path,
                                character_definitions,
                                scene_characters
                        ):
                            image_paths.append(image_path)
                            # Add this panel to context tracker
                            context_tracker.add_panel(idx, scene, scene_characters, character_descriptions)
                            logging.info(f"Consistent image generated for scene {idx}")
                            attempt_successful = True
                            break
                        else:
                            logging.warning(f"Inconsistency detected, retrying... ({attempt + 1}/{max_attempts})")
                            continue

                    if image_url:
                        image_response = requests.get(image_url)
                        if image_response.status_code == 200:
                            with open(image_path, 'wb') as img_file:
                                img_file.write(image_response.content)

                            if verify_character_consistency_enhanced(
                                    image_path,
                                    character_definitions,
                                    scene_characters
                            ):
                                image_paths.append(image_path)
                                # Add this panel to context tracker
                                context_tracker.add_panel(idx, scene, scene_characters, character_descriptions)
                                logging.info(f"Consistent image generated for scene {idx}")
                                attempt_successful = True
                                break
                            else:
                                logging.warning(f"Inconsistency detected, retrying... ({attempt + 1}/{max_attempts})")

            except Exception as e:
                logging.error(f"Error generating image for Scene {idx}, attempt {attempt + 1}: {e}")

        # Use fallback if all attempts fail
        if not attempt_successful:
            if fallback_image_generation(prompt, image_path):
                image_paths.append(image_path)
                # Still add to context tracker even for fallback
                context_tracker.add_panel(idx, scene, scene_characters, character_descriptions)
                logging.info(f"Fallback image saved for scene {idx}")

    return image_paths


# -------------------------------------------------------
# 15) BASIC IMAGE GENERATION (WITHOUT CONSISTENCY) - COLOR
# -------------------------------------------------------
def generate_comic_images(story_text, title, story_dir):
    """Basic color comic image generation without consistency checking."""
    logging.info("Starting basic color comic image generation.")

    if not openai_client:
        logging.error("OpenAI client not initialized. Cannot generate images.")
        return []

    image_dir = os.path.join(story_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    scenes = extract_key_elements(story_text)
    image_paths = []

    for idx, scene in enumerate(scenes, start=1):
        image_path = os.path.join(image_dir, f"scene_{idx}.png")
        prompt = f"Full color comic style art with vibrant colors of the following scene: {scene}"

        try:
            response = openai_client.images.generate(
                model="gpt-image-1",
                prompt=prompt,
                size="1024x1024",
                quality="high",
                n=1,
            )

            if hasattr(response, 'data') and response.data and len(response.data) > 0:
                data_item = response.data[0]
                image_url = None

                if hasattr(data_item, 'url') and data_item.url:
                    image_url = data_item.url
                elif hasattr(data_item, 'b64_json') and data_item.b64_json:
                    img_data = base64.b64decode(data_item.b64_json)
                    with open(image_path, 'wb') as img_file:
                        img_file.write(img_data)
                    image_paths.append(image_path)
                    continue

            if image_url:
                image_response = requests.get(image_url)
                if image_response.status_code == 200:
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_response.content)
                    image_paths.append(image_path)

        except Exception as e:
            logging.error(f"Error generating image for Scene {idx}: {e}")
            if fallback_image_generation(prompt, image_path):
                image_paths.append(image_path)

    return image_paths


# -------------------------------------------------------
# 16) AUDIO GENERATION WITH EMOTION
# -------------------------------------------------------
def analyze_scene_emotion(scene_text):
    """Analyzes the emotional tone of a scene."""
    voice_name = "nova"

    emotion_keywords = {
        "excited": ["excited", "thrilled", "ecstatic", "joy", "happy", "elated", "enthusiastic"],
        "sad": ["sad", "melancholy", "depressed", "grief", "sorrow", "miserable", "unhappy"],
        "angry": ["angry", "furious", "rage", "irate", "livid", "outraged", "mad"],
        "fearful": ["fear", "scared", "terrified", "frightened", "panicked", "anxious", "dread"],
        "calm": ["calm", "peaceful", "serene", "tranquil", "relaxed", "gentle", "soothing"]
    }

    scene_lower = scene_text.lower()
    emotion_counts = {emotion: 0 for emotion in emotion_keywords}

    for emotion, keywords in emotion_keywords.items():
        for keyword in keywords:
            emotion_counts[emotion] += scene_lower.count(keyword)

    dominant_emotion = max(emotion_counts.items(), key=lambda x: x[1])[0] if any(emotion_counts.values()) else None

    if dominant_emotion == "excited" or dominant_emotion == "angry":
        voice_name = "shimmer"
    elif dominant_emotion == "sad" or dominant_emotion == "fearful":
        voice_name = "onyx"
    elif dominant_emotion == "calm":
        voice_name = "nova"

    has_exclamations = "!" in scene_text
    has_dialogue = '"' in scene_text or "'" in scene_text

    if has_dialogue and (scene_text.count('"') > 5 or scene_text.count("'") > 5):
        voice_name = "echo"

    if has_exclamations and scene_text.count("!") > 3:
        voice_name = "shimmer"

    return voice_name


def enhance_text_for_tts(scene_text):
    """Enhances the text to improve emotional expression in TTS."""
    enhanced_text = scene_text

    enhanced_text = enhanced_text.replace(". ", ". <break time='0.3s'> ")

    emphasis_words = ["suddenly", "never", "always", "terrified", "ecstatic", "furious", "absolutely"]
    for word in emphasis_words:
        pattern = r"\b" + word + r"\b"
        replacement = word.upper()
        enhanced_text = re.sub(pattern, replacement, enhanced_text, flags=re.IGNORECASE)

    if '"' in enhanced_text:
        enhanced_text = enhanced_text.replace(' "', ' <break time="0.2s"> "')

    return enhanced_text


def generate_audio_narration(story_text, title, output_dir="storyboard_assets"):
    """Generates emotionally expressive audio narration."""
    logging.info("Starting audio narration generation.")

    if not openai_client:
        logging.error("OpenAI client not initialized. Cannot generate audio narration.")
        return [], output_dir

    safe_title = "".join(c if c.isalnum() else "_" for c in title)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    story_dir = os.path.join(output_dir, f"{safe_title}_{timestamp}")
    audio_dir = os.path.join(story_dir, "audio")
    os.makedirs(audio_dir, exist_ok=True)

    scenes = extract_key_elements(story_text)
    audio_paths = []

    for idx, scene in enumerate(scenes, start=1):
        audio_path = os.path.join(audio_dir, f"scene_{idx}.mp3")
        try:
            voice = analyze_scene_emotion(scene)
            enhanced_scene = enhance_text_for_tts(scene)

            response = openai_client.audio.speech.create(
                model="tts-1-hd",
                voice=voice,
                input=enhanced_scene,
                speed=0.95
            )

            response.stream_to_file(audio_path)
            audio_paths.append(audio_path)
            logging.info(f"Audio saved with {voice} voice: {audio_path}")
        except Exception as e:
            logging.error(f"Error generating audio for Scene {idx}: {e}")

    return audio_paths, story_dir


# -------------------------------------------------------
# 17) PDF CREATION
# -------------------------------------------------------
class UnicodeAwarePDF(FPDF):
    """Extended FPDF class with Unicode support."""

    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        try:
            self.add_font('DejaVu', '', 'DejaVuSansCondensed.ttf', uni=True)
            self.add_font('DejaVu', 'B', 'DejaVuSansCondensed-Bold.ttf', uni=True)
        except:
            pass


def create_storyboard_with_audio(image_paths, audio_paths, story_text, title, story_dir):
    """Creates a PDF storyboard and PowerPoint presentation."""
    logging.info("Creating storyboard PDF and PowerPoint.")

    pdf_path = os.path.join(story_dir, f"{title}_storyboard.pdf")
    pptx_path = create_storyboard_pptx(image_paths, audio_paths, story_text, title, story_dir)

    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(190, 10, sanitize_text_for_pdf(f"Storyboard: {title}"), ln=True, align="C")
        pdf.ln(10)

        scenes = extract_key_elements(story_text)

        for idx, scene in enumerate(scenes, start=1):
            pdf.add_page()
            pdf.set_font('Arial', 'B', 14)
            pdf.cell(190, 10, f"Scene {idx}", ln=True)

            if idx <= len(image_paths) and os.path.exists(image_paths[idx - 1]):
                pdf.image(image_paths[idx - 1], x=10, y=30, w=190)
                y_position = 150
            else:
                y_position = 40

            pdf.set_xy(10, y_position)
            pdf.set_font('Arial', '', 10)
            sanitized_scene = sanitize_text_for_pdf(scene)
            pdf.multi_cell(190, 5, sanitized_scene)

            if idx <= len(audio_paths) and os.path.exists(audio_paths[idx - 1]):
                pdf.set_xy(10, 280)
                pdf.set_font('Arial', 'I', 8)
                pdf.cell(190, 10, f"Audio: {os.path.basename(audio_paths[idx - 1])}", ln=True)

        pdf.output(pdf_path)
        logging.info(f"PDF created: {pdf_path}")
        return pdf_path, pptx_path
    except Exception as e:
        logging.error(f"Error creating PDF: {e}")
        return None, pptx_path


# -------------------------------------------------------
# 18) POWERPOINT CREATION
# -------------------------------------------------------
def create_storyboard_pptx(image_paths, audio_paths, story_text, title, story_dir):
    """Creates a PowerPoint presentation."""
    logging.info("Creating storyboard PowerPoint presentation.")

    pptx_path = os.path.join(story_dir, f"{title}_storyboard.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)

    title_shape = title_slide.shapes.title
    title_shape.text = f"Storyboard: {title}"
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(79, 70, 229)

    scenes = extract_key_elements(story_text)

    for idx, scene in enumerate(scenes, start=1):
        content_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(content_slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = f"Scene {idx}"
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        if idx <= len(image_paths) and os.path.exists(image_paths[idx - 1]):
            img_width = Inches(8)
            img_left = (prs.slide_width - img_width) / 2
            img_top = Inches(1.5)

            try:
                slide.shapes.add_picture(image_paths[idx - 1], img_left, img_top, width=img_width)
            except Exception as e:
                logging.error(f"Error adding image to slide {idx}: {e}")

        text_left = Inches(1)
        text_top = Inches(5.5)
        text_width = Inches(11.33)
        text_height = Inches(1.5)

        txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = scene[:500] + "..." if len(scene) > 500 else scene
        p.font.size = Pt(12)

    try:
        prs.save(pptx_path)
        logging.info(f"PowerPoint created: {pptx_path}")
        return pptx_path
    except Exception as e:
        logging.error(f"Error creating PowerPoint: {e}")
        return None


# -------------------------------------------------------
# 19) DATABASE OPERATIONS
# -------------------------------------------------------
def save_to_mysql(title, original_story, enhanced_story, audio_files, image_files, pdf_path, pptx_path, user_id=None):
    """Saves story data to MySQL."""
    logging.info("Saving story data to MySQL...")

    audio_json = json.dumps(audio_files)
    image_json = json.dumps(image_files)

    if isinstance(pdf_path, tuple):
        pdf_local = os.path.abspath(pdf_path[0]) if pdf_path[0] and os.path.exists(pdf_path[0]) else ""
    else:
        pdf_local = os.path.abspath(pdf_path) if pdf_path and os.path.exists(pdf_path) else ""

    if isinstance(pptx_path, tuple):
        pptx_local = os.path.abspath(pptx_path[1]) if len(pptx_path) > 1 and pptx_path[1] and os.path.exists(
            pptx_path[1]) else ""
    else:
        pptx_local = os.path.abspath(pptx_path) if pptx_path and os.path.exists(pptx_path) else ""

    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor()

        if user_id:
            insert_query = """
                INSERT INTO story_data 
                (title, story_outline, enhanced_story, audio_urls, image_urls, pdf_url, pptx_url, user_id)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
            """
            cursor.execute(
                insert_query,
                (title, original_story, enhanced_story, audio_json, image_json, pdf_local, pptx_local, user_id)
            )
        else:
            insert_query = """
                INSERT INTO story_data 
                (title, story_outline, enhanced_story, audio_urls, image_urls, pdf_url, pptx_url)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
            """
            cursor.execute(
                insert_query,
                (title, original_story, enhanced_story, audio_json, image_json, pdf_local, pptx_local)
            )

        conn.commit()
        record_id = cursor.lastrowid
        cursor.close()
        conn.close()
        logging.info(f"MySQL insert successful with ID: {record_id}")
        return record_id

    except Exception as e:
        logging.error(f"Error inserting data into MySQL: {e}")
        return None


def save_story_with_character_styles(title, original_story, enhanced_story, audio_files, image_files, pdf_path,
                                     pptx_path, character_styles, user_id=None):
    """Saves story data with character styles to MySQL."""
    logging.info("Saving story data with character styles to MySQL...")

    audio_json = json.dumps(audio_files)
    image_json = json.dumps(image_files)
    character_styles_json = json.dumps(character_styles) if character_styles else None

    if isinstance(pdf_path, tuple):
        pdf_local = os.path.abspath(pdf_path[0]) if pdf_path[0] and os.path.exists(pdf_path[0]) else ""
    else:
        pdf_local = os.path.abspath(pdf_path) if pdf_path and os.path.exists(pdf_path) else ""

    if isinstance(pptx_path, tuple):
        pptx_local = os.path.abspath(pptx_path[1]) if len(pptx_path) > 1 and pptx_path[1] and os.path.exists(
            pptx_path[1]) else ""
    else:
        pptx_local = os.path.abspath(pptx_path) if pptx_path and os.path.exists(pptx_path) else ""

    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor()

        if user_id:
            insert_query = """
                INSERT INTO story_data 
                (title, story_outline, enhanced_story, audio_urls, image_urls, pdf_url, pptx_url, character_styles, user_id)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
            """
            cursor.execute(
                insert_query,
                (title, original_story, enhanced_story, audio_json, image_json, pdf_local, pptx_local,
                 character_styles_json, user_id)
            )
        else:
            insert_query = """
                INSERT INTO story_data 
                (title, story_outline, enhanced_story, audio_urls, image_urls, pdf_url, pptx_url, character_styles)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
            """
            cursor.execute(
                insert_query,
                (title, original_story, enhanced_story, audio_json, image_json, pdf_local, pptx_local,
                 character_styles_json)
            )

        conn.commit()
        record_id = cursor.lastrowid
        cursor.close()
        conn.close()
        logging.info(f"MySQL insert successful with ID: {record_id}")
        return record_id

    except Exception as e:
        logging.error(f"Error inserting data into MySQL: {e}")
        return None


# -------------------------------------------------------
# 20) RETRIEVAL FUNCTIONS
# -------------------------------------------------------
def get_story_history(limit=10, user_id=None):
    """Retrieves the most recent stories from the database."""
    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor(pymysql.cursors.DictCursor)

        if user_id:
            query = """
                SELECT id, title, SUBSTRING(story_outline, 1, 100) AS preview, 
                       created_at, pdf_url, pptx_url, user_id
                FROM story_data 
                WHERE user_id = %s
                ORDER BY created_at DESC
                LIMIT %s
            """
            cursor.execute(query, (user_id, limit))
        else:
            query = """
                SELECT id, title, SUBSTRING(story_outline, 1, 100) AS preview, 
                       created_at, pdf_url, pptx_url, user_id
                FROM story_data 
                ORDER BY created_at DESC
                LIMIT %s
            """
            cursor.execute(query, (limit,))

        stories = cursor.fetchall()
        cursor.close()
        conn.close()
        return stories

    except Exception as e:
        logging.error(f"Error retrieving story history: {e}")
        return []


def get_story_details(story_id):
    """Retrieves full details of a specific story."""
    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor(pymysql.cursors.DictCursor)

        query = "SELECT * FROM story_data WHERE id = %s"
        cursor.execute(query, (story_id,))
        story = cursor.fetchone()

        cursor.close()
        conn.close()

        if story:
            story['audio_urls'] = json.loads(story['audio_urls'])
            story['image_urls'] = json.loads(story['image_urls'])

            if 'character_styles' in story and story['character_styles']:
                try:
                    story['character_styles'] = json.loads(story['character_styles'])
                except:
                    story['character_styles'] = {}
            else:
                story['character_styles'] = {}

        return story

    except Exception as e:
        logging.error(f"Error retrieving story details: {e}")
        return None


# -------------------------------------------------------
# 21) UPDATE FUNCTIONS
# -------------------------------------------------------
def update_story_data(story_id, title, story_outline, enhanced_story, audio_files, image_files, pdf_path, pptx_path,
                      character_styles=None, user_id=None):
    """Updates an existing story in the database."""
    logging.info(f"Updating story ID {story_id} with new data")

    audio_json = json.dumps(audio_files)
    image_json = json.dumps(image_files)
    character_styles_json = json.dumps(character_styles) if character_styles else None

    if isinstance(pdf_path, tuple):
        pdf_local = os.path.abspath(pdf_path[0]) if pdf_path[0] and os.path.exists(pdf_path[0]) else ""
    else:
        pdf_local = os.path.abspath(pdf_path) if pdf_path and os.path.exists(pdf_path) else ""

    if isinstance(pptx_path, tuple):
        pptx_local = os.path.abspath(pptx_path[1]) if len(pptx_path) > 1 and pptx_path[1] and os.path.exists(
            pptx_path[1]) else ""
    else:
        pptx_local = os.path.abspath(pptx_path) if pptx_path and os.path.exists(pptx_path) else ""

    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor()

        if user_id:
            update_query = """
                UPDATE story_data 
                SET title = %s, 
                    story_outline = COALESCE(%s, story_outline), 
                    enhanced_story = COALESCE(%s, enhanced_story), 
                    audio_urls = %s, 
                    image_urls = %s, 
                    pdf_url = %s,
                    pptx_url = %s,
                    character_styles = %s
                WHERE id = %s AND user_id = %s
            """
            cursor.execute(
                update_query,
                (title, story_outline, enhanced_story, audio_json, image_json,
                 pdf_local, pptx_local, character_styles_json, story_id, user_id)
            )
        else:
            update_query = """
                UPDATE story_data 
                SET title = %s, 
                    story_outline = COALESCE(%s, story_outline), 
                    enhanced_story = COALESCE(%s, enhanced_story), 
                    audio_urls = %s, 
                    image_urls = %s, 
                    pdf_url = %s,
                    pptx_url = %s,
                    character_styles = %s
                WHERE id = %s
            """
            cursor.execute(
                update_query,
                (title, story_outline, enhanced_story, audio_json, image_json,
                 pdf_local, pptx_local, character_styles_json, story_id)
            )

        rows_affected = cursor.rowcount
        conn.commit()
        cursor.close()
        conn.close()

        if rows_affected > 0:
            logging.info(f"Successfully updated story ID {story_id}")
            return True
        else:
            logging.warning(f"No rows affected when updating story ID {story_id}")
            return False

    except Exception as e:
        logging.error(f"Error updating story data: {e}")
        return False


def recreate_story_assets(story_id, story_text, title, character_styles=None, consistency_level=0.8):
    """Regenerates all assets for an existing story."""
    logging.info(f"Regenerating assets for story ID {story_id}")

    try:
        safe_title = "".join(c if c.isalnum() else "_" for c in title)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        story_dir = os.path.join("storyboard_assets", f"{safe_title}_{timestamp}")
        os.makedirs(story_dir, exist_ok=True)

        audio_files, _ = generate_audio_narration(story_text, title, story_dir)

        if character_styles and len(character_styles) > 0:
            image_files = generate_comic_images_with_consistency(
                story_text, title, story_dir, character_styles, consistency_level
            )
        else:
            image_files = generate_comic_images(story_text, title, story_dir)

        pdf_path, pptx_path = create_storyboard_with_audio(image_files, audio_files, story_text, title, story_dir)

        return {
            "audio_files": audio_files,
            "image_files": image_files,
            "pdf_path": pdf_path,
            "pptx_path": pptx_path,
            "story_dir": story_dir
        }

    except Exception as e:
        logging.error(f"Error regenerating story assets: {e}")
        return None


# -------------------------------------------------------
# 22) UTILITY FUNCTIONS
# -------------------------------------------------------
def get_user_stories(user_id, limit=10):
    """Retrieves stories for a specific user."""
    return get_story_history(limit, user_id)


def get_user_by_id(user_id):
    """Retrieves user information by ID."""
    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor(pymysql.cursors.DictCursor)

        query = "SELECT id, username, email, created_at FROM users WHERE id = %s"
        cursor.execute(query, (user_id,))
        user = cursor.fetchone()

        cursor.close()
        conn.close()
        return user

    except Exception as e:
        logging.error(f"Error retrieving user information: {e}")
        return None


def save_character_styles(story_id, character_styles):
    """Save character styles for a story."""
    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor()

        character_styles_json = json.dumps(character_styles)

        cursor.execute(
            "UPDATE story_data SET character_styles = %s WHERE id = %s",
            (character_styles_json, story_id)
        )

        conn.commit()
        cursor.close()
        conn.close()
        return True

    except Exception as e:
        logging.error(f"Error saving character styles: {e}")
        return False


def get_character_styles(story_id):
    """Get character styles for a story."""
    try:
        conn = pymysql.connect(
            host=mysql_host, user=mysql_user, password=mysql_password, database=mysql_db
        )
        cursor = conn.cursor(pymysql.cursors.DictCursor)

        cursor.execute("SELECT character_styles FROM story_data WHERE id = %s", (story_id,))
        result = cursor.fetchone()

        cursor.close()
        conn.close()

        if result and result['character_styles']:
            return json.loads(result['character_styles'])

        return None

    except Exception as e:
        logging.error(f"Error retrieving character styles: {e}")
        return None


def analyze_story_for_character_consistency(story_text, character_styles=None):
    """Analyze story for character consistency issues."""
    logging.info("Analyzing story for character consistency...")

    if not openai_client:
        return {"issues": [], "summary": "Analysis failed due to missing OpenAI client."}

    if not character_styles:
        character_styles = extract_characters_from_story(story_text)

    try:
        messages = [
            {"role": "system", "content": "You are a literary analyst specializing in character consistency."},
            {"role": "user", "content": f"""
                Analyze the following story for character consistency issues. 
                Check for:
                1. Physical appearance inconsistencies
                2. Personality/behavior inconsistencies
                3. Timeline/location inconsistencies related to characters

                Character information:
                {json.dumps(character_styles, indent=2)}

                Story:
                {story_text}

                Format your response as a JSON object with:
                - "issues": Array of consistency issues, each with "type", "description", "severity" (1-5)
                - "summary": Brief overall assessment of character consistency
            """}
        ]

        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            temperature=0.3,
            max_tokens=1500,
            response_format={"type": "json_object"}
        )

        analysis = json.loads(response.choices[0].message.content)
        return analysis

    except Exception as e:
        logging.error(f"Error analyzing character consistency: {e}")
        return {"issues": [], "summary": "Analysis failed due to an error."}


def generate_consistent_image_prompt(scene_text, character_styles, consistency_level=0.8, panel_context=None):
    """Generate an image prompt with character sheet references and previous panel context."""
    logging.info("Generating style-consistent image prompt with context...")

    try:
        # Create character sheet descriptions
        character_sheets = []
        for char_name, char_info in character_styles.items():
            sheet_parts = []

            if 'height' in char_info:
                sheet_parts.append(f"Height: {char_info['height']}")
            if 'age' in char_info:
                sheet_parts.append(f"Age: {char_info['age']}")
            if 'face_shape' in char_info:
                sheet_parts.append(f"Face: {char_info['face_shape']}")
            if 'hair' in char_info:
                sheet_parts.append(f"Hair: {char_info['hair']}")
            if 'eyes' in char_info:
                sheet_parts.append(f"Eyes: {char_info['eyes']}")
            if 'skin_tone' in char_info:
                sheet_parts.append(f"Skin: {char_info['skin_tone']}")
            if 'build' in char_info:
                sheet_parts.append(f"Build: {char_info['build']}")
            if 'clothing' in char_info:
                sheet_parts.append(f"Clothing: {char_info['clothing']}")
            if 'distinctive_features' in char_info:
                sheet_parts.append(f"Distinctive: {char_info['distinctive_features']}")
            if 'typical_pose' in char_info:
                sheet_parts.append(f"Pose/Expression: {char_info['typical_pose']}")

            sheet = f"CHARACTER SHEET - {char_name}:\n" + "\n".join(f"- {part}" for part in sheet_parts)
            character_sheets.append(sheet)

        sheets_text = "\n\n".join(character_sheets)

        # Add context from previous panels if available
        context_text = panel_context if panel_context else ""

        # Create the structured prompt with COLOR and CONTEXT
        enhanced_prompt = f"""Create a FULL COLOR comic style illustration with vibrant colors.

CRITICAL CONSISTENCY REQUIREMENTS:
- ALL characters MUST match their character sheets EXACTLY
- Maintain precise colors, proportions, features, and clothing
- Use the character sheets as absolute reference
- Characters must be instantly recognizable from previous panels
- Consistent color palette throughout all panels

{context_text}

{sheets_text}

SCENE TO ILLUSTRATE:
{scene_text}

STYLE NOTES:
- FULL COLOR comic art style with rich, vibrant colors
- Clean line art with professional coloring
- Consistent character appearances and colors as defined in sheets
- Professional comic book quality with dynamic shading
- Use appropriate atmospheric colors for mood and setting
- Maintain continuity with previous panels as described in context"""

        return enhanced_prompt

    except Exception as e:
        logging.error(f"Error generating consistent image prompt: {e}")
        return f"Full color comic style art with vibrant colors of the following scene: {scene_text}"


def extract_detailed_character_appearance(character_name, character_styles, previous_panel_details=None):
    """Extract detailed character appearance for context tracking."""
    appearance_details = []

    if character_name in character_styles:
        char_info = character_styles[character_name]

        # Base details from character sheet
        if 'hair' in char_info:
            appearance_details.append(f"hair: {char_info['hair']}")
        if 'eyes' in char_info:
            appearance_details.append(f"eyes: {char_info['eyes']}")
        if 'clothing' in char_info:
            appearance_details.append(f"wearing: {char_info['clothing']}")
        if 'distinctive_features' in char_info:
            appearance_details.append(f"features: {char_info['distinctive_features']}")

        # If we have previous panel details, add any dynamic information
        if previous_panel_details:
            if 'pose' in previous_panel_details:
                appearance_details.append(f"last pose: {previous_panel_details['pose']}")
            if 'expression' in previous_panel_details:
                appearance_details.append(f"last expression: {previous_panel_details['expression']}")

    return ", ".join(appearance_details)